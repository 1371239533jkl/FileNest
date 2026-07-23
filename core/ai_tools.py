"""
AI 工具注册与执行系统 —— 让 LLM 自主调度多种工具完成复杂任务。

支持的预设工具：
- search_files  : 检索本地文件数据库（保留原有能力）
- search_content: 检索已建立索引的本地文档正文，并返回来源片段
- search_web    : 联网搜索，获取实时信息
- read_file     : 读取本地文件内容

用法：
    registry = ToolRegistry()
    registry.register(search_files_tool)
    schemas = registry.get_tool_schemas()  # 传给 LLM 的 tools 参数
    result = registry.execute("search_web", {"query": "Python 3.13"})
"""

import json
import os
import re
from pathlib import Path
from typing import Optional, Callable, Any
from dataclasses import dataclass, field

import httpx

from utils.logger import logger


# ══════════════════════════════════════════════════════════════════════════════
# 工具定义
# ══════════════════════════════════════════════════════════════════════════════

@dataclass
class ToolDefinition:
    """单个工具定义"""
    name: str                              # 工具名，如 "search_files"
    description: str                       # LLM 可读的功能描述
    parameters: dict                       # JSON Schema 格式的参数定义
    handler: Callable[..., str]            # 执行函数，接收 **kwargs，返回文本结果
    requires_db: bool = False              # 是否需要数据库连接
    db_manager: Any = None                 # 数据库管理器实例（延迟注入）

    def execute(self, arguments: dict) -> str:
        """执行工具并返回文本结果"""
        try:
            result = self.handler(**(arguments or {}))
            # 截断过长结果，避免撑爆上下文
            if isinstance(result, str) and len(result) > 8000:
                result = result[:8000] + "\n...(结果已截断)"
            return result
        except Exception as e:
            logger.error(f"工具 {self.name} 执行失败: {e}")
            return f"[工具执行错误] {self.name}: {e}"


# ══════════════════════════════════════════════════════════════════════════════
# 工具注册表
# ══════════════════════════════════════════════════════════════════════════════

class ToolRegistry:
    """工具注册与调度中心"""

    def __init__(self):
        self._tools: dict[str, ToolDefinition] = {}

    def register(self, tool: ToolDefinition) -> None:
        """注册一个工具"""
        self._tools[tool.name] = tool
        logger.debug(f"工具已注册: {tool.name}")

    def unregister(self, name: str) -> None:
        """注销一个工具"""
        self._tools.pop(name, None)

    def get(self, name: str) -> Optional[ToolDefinition]:
        """获取指定工具"""
        return self._tools.get(name)

    def list_tools(self) -> list[str]:
        """列出所有已注册工具名"""
        return list(self._tools.keys())

    def get_tool_schemas(self) -> list[dict]:
        """生成 OpenAI 兼容的 tools 参数列表（仅描述，不含 handler）"""
        return [
            {
                "type": "function",
                "function": {
                    "name": t.name,
                    "description": t.description,
                    "parameters": t.parameters,
                }
            }
            for t in self._tools.values()
        ]

    def execute(self, tool_name: str, arguments: dict) -> str:
        """执行指定工具并返回文本结果"""
        tool = self._tools.get(tool_name)
        if not tool:
            return f"[错误] 未知工具: {tool_name}，可用工具: {', '.join(self.list_tools())}"
        return tool.execute(arguments)


# ══════════════════════════════════════════════════════════════════════════════
# 预设工具：文件搜索
# ══════════════════════════════════════════════════════════════════════════════

def _search_files_handler(query: str = None, file_type: str = None,
                          start_date: str = None, end_date: str = None,
                          min_size: int = None, max_size: int = None,
                          max_results: int = 500, _db=None, _ai_layer=None) -> str:
    """搜索本地文件数据库 —— 支持关键词、日期范围、大小范围、文件类型过滤"""
    if _db is None:
        return "[错误] 文件搜索工具未连接数据库"

    try:
        # 使用 FileDAO.search() 的结构化查询
        from database.models import FileDAO
        from utils.display_utils import format_size as fmt

        file_dao = FileDAO(_db)

        # 处理 query：如果不是纯文件名关键词，可能包含日期/大小语义
        # FileDAO.search 只接受单个 name，所以取 query 整体或拆第一个词
        name = query.strip() if query else None

        rows = file_dao.search(
            name=name,
            file_type=file_type,
            min_size=min_size,
            max_size=max_size,
            start_date=start_date,
            end_date=end_date,
        )

        # 限制返回数量
        rows = rows[:max_results]

        if not rows:
            parts = []
            if query:
                parts.append(f"关键词: {query}")
            if file_type:
                from config import FILE_TYPE_NAMES
                parts.append(f"类型: {FILE_TYPE_NAMES.get(file_type, file_type)}")
            if start_date or end_date:
                parts.append(f"日期: {start_date or '不限'} ~ {end_date or '不限'}")
            if min_size is not None or max_size is not None:
                sz = []
                if min_size is not None:
                    sz.append(f"≥{_format_bytes(min_size)}")
                if max_size is not None:
                    sz.append(f"≤{_format_bytes(max_size)}")
                parts.append(f"大小: {' '.join(sz)}")
            desc = "、".join(parts) if parts else "该条件"
            return f"未找到符合条件（{desc}）的文件。"

        # 按修改时间倒序
        from datetime import datetime
        def _sort_key(r):
            mt = r.get('modify_time')
            if isinstance(mt, datetime):
                return mt
            if isinstance(mt, str):
                try:
                    return datetime.strptime(mt[:19], "%Y-%m-%d %H:%M:%S")
                except Exception:
                    pass
            return datetime.min
        rows.sort(key=_sort_key, reverse=True)

        total_size = sum(r.get('file_size', 0) for r in rows)
        lines = [f"找到 {len(rows)} 个文件（总大小: {_format_bytes(total_size)}）:"]
        for r in rows:
            size_str = _format_bytes(r.get('file_size', 0))
            mt = r.get('modify_time', '')
            if isinstance(mt, datetime):
                mt = mt.strftime("%Y-%m-%d")
            elif isinstance(mt, str) and len(mt) >= 10:
                mt = mt[:10]
            lines.append(
                f"- {r['file_name']} | {size_str} | 修改:{mt} | "
                f"类型:{r.get('file_type','?')} | {r['file_path']}"
            )
        return "\n".join(lines)

    except Exception as e:
        return f"[文件搜索错误] {e}"


_search_files_schema = {
    "type": "object",
    "properties": {
        "query": {
            "type": "string",
            "description": "文件名或路径关键词。如果用户说'近一年的文档'且没有具体文件名，query可以不传，用start_date+file_type替代"
        },
        "file_type": {
            "type": "string",
            "enum": ["image", "document", "code", "video", "audio", "archive", "executable", "font", "other"],
            "description": "文件类型过滤，如用户说'文档'则传 'document'"
        },
        "start_date": {
            "type": "string",
            "description": "文件修改时间起始日期，格式 YYYY-MM-DD。用户说'近一年'推算为一年前今天"
        },
        "end_date": {
            "type": "string",
            "description": "文件修改时间截止日期，格式 YYYY-MM-DD。用户说'近一年'则end_date为今天"
        },
        "min_size": {
            "type": "integer",
            "description": "最小文件大小（字节）。用户说'大于100MB'则传 104857600"
        },
        "max_size": {
            "type": "integer",
            "description": "最大文件大小（字节）。用户说'小于1GB'则传 1073741824"
        },
        "max_results": {
            "type": "integer",
            "description": "最大返回条数。默认 500，结果少时传小值节省上下文，需要大量结果时才传大值（最多 2000）。",
            "default": 500
        }
    },
    "required": []
}


def create_search_files_tool(db_manager=None, ai_layer=None) -> ToolDefinition:
    """创建文件搜索工具（注入数据库依赖）"""
    def handler(**kwargs):
        return _search_files_handler(**kwargs, _db=db_manager, _ai_layer=ai_layer)
    return ToolDefinition(
        name="search_files",
        description="搜索本地文件数据库。当你需要查找用户电脑上的文件时使用此工具。支持按文件名、路径关键词搜索，可按文件类型过滤。",
        parameters=_search_files_schema,
        handler=handler,
        requires_db=True,
        db_manager=db_manager,
    )


# ══════════════════════════════════════════════════════════════════════════════
# 预设工具：已索引正文检索
# ══════════════════════════════════════════════════════════════════════════════

def _search_content_handler(query: str, max_results: int = 8, _db=None) -> str:
    """Search only the local extracted-text index and return cited snippets."""
    if _db is None:
        return "[错误] 正文检索工具未连接数据库"
    query = (query or '').strip()
    if not query:
        return "[错误] 请提供需要检索的正文关键词"

    try:
        from database.models import FileContentDAO

        limit = max(1, min(int(max_results or 8), 20))
        rows = FileContentDAO(_db).search(query, limit=limit)
        if not rows:
            return ("未在已建立正文索引的文件中找到匹配内容。"
                    "可提醒用户先在搜索页或扫描页执行“构建正文索引”。")

        lines = [f"正文检索“{query}”命中 {len(rows)} 个文件。回答时请引用 [来源 N]："]
        for index, row in enumerate(rows, 1):
            snippet = (row.get('content_snippet') or '').replace('\n', ' ').strip()
            lines.append(
                f"[来源 {index}] {row.get('file_name', '未知文件')} | "
                f"路径: {row.get('file_path', '')} | 摘录: {snippet}"
            )
        return "\n".join(lines)
    except (TypeError, ValueError) as exc:
        return f"[正文检索错误] 参数无效: {exc}"
    except Exception as exc:
        logger.error(f"正文检索失败: {exc}")
        return f"[正文检索错误] {exc}"


_search_content_schema = {
    "type": "object",
    "properties": {
        "query": {
            "type": "string",
            "description": "要在已索引文档正文中查找的关键词或短语。"
        },
        "max_results": {
            "type": "integer",
            "description": "最多返回的来源数量，默认 8，最大 20。",
            "default": 8,
        },
    },
    "required": ["query"],
}


def create_search_content_tool(db_manager=None) -> ToolDefinition:
    """Create a read-only local content-search tool with source citations."""
    return ToolDefinition(
        name="search_content",
        description=("检索用户已建立索引的本地文档正文。适合回答文档中提到什么、"
                     "查找某个术语或根据文档内容归纳时使用。结果包含可追溯的文件来源和摘录。"),
        parameters=_search_content_schema,
        handler=lambda **kwargs: _search_content_handler(**kwargs, _db=db_manager),
        requires_db=True,
        db_manager=db_manager,
    )


# ══════════════════════════════════════════════════════════════════════════════
# 预设工具：联网搜索
# ══════════════════════════════════════════════════════════════════════════════

_WEB_SEARCH_CLIENT: Optional[httpx.Client] = None


def _get_web_client() -> httpx.Client:
    global _WEB_SEARCH_CLIENT
    if _WEB_SEARCH_CLIENT is None:
        _WEB_SEARCH_CLIENT = httpx.Client(
            timeout=httpx.Timeout(30.0),
            headers={"User-Agent": "SmartFileManager/2.0 (AI Assistant)"},
            follow_redirects=True,
        )
    return _WEB_SEARCH_CLIENT


# SearXNG 公共实例列表（免费、开源、无需 API Key）
_SEARXNG_INSTANCES = [
    "https://search.sapti.me",
    "https://searx.be",
    "https://search.bus-hit.me",
    "https://searx.tiekoetter.com",
    "https://search.rhscz.eu",
]


def _get_bocha_key() -> Optional[str]:
    """获取博查 API Key（环境变量或 config）"""
    key = os.getenv("BOCHA_API_KEY", "")
    if key:
        return key
    try:
        from config import AI_CONFIG
        return AI_CONFIG.get("bocha_api_key", "") or ""
    except Exception:
        return ""


def _search_with_bocha(query: str, max_results: int) -> Optional[str]:
    """博查 Web Search API（免费 2000次/月，需 API Key）
    注册地址: https://open.bochaai.com
    """
    api_key = _get_bocha_key()
    if not api_key:
        logger.debug("博查 API Key 未配置，跳过")
        return None
    try:
        with httpx.Client(
            timeout=httpx.Timeout(8.0),
            headers={
                "Authorization": f"Bearer {api_key}",
                "Content-Type": "application/json",
                "User-Agent": "SmartFileManager/2.0",
            },
        ) as client:
            logger.debug(f"博查搜索: {query[:30]}...")
            resp = client.post(
                "https://api.bocha.cn/v1/web-search",
                json={"query": query, "freshness": "noLimit", "summary": True, "count": max_results},
            )
            logger.debug(f"博查响应: HTTP {resp.status_code}")
            if resp.status_code != 200:
                logger.warning(f"博查 API 错误: HTTP {resp.status_code}, {resp.text[:200]}")
                return None
            data = resp.json()
            if data.get("code") != 200:
                logger.warning(f"博查业务错误: code={data.get('code')}, msg={data.get('msg', 'unknown')}")
                return None
            entries = data.get("data", {}).get("webPages", {}).get("value", [])
            if not entries:
                logger.debug("博查返回空结果")
                return None
            lines = [f"网页搜索 '{query}' 的结果 (via Bocha):"]
            for i, e in enumerate(entries[:max_results]):
                lines.append(f"{i+1}. {e.get('name','无标题')}")
                snippet = e.get("snippet", e.get("summary", ""))
                if snippet:
                    lines.append(f"   {snippet[:200]}")
                if e.get("url"):
                    lines.append(f"   链接: {e['url']}")
            logger.info(f"博查搜索成功: {len(entries)} 条结果")
            return "\n".join(lines)
    except httpx.TimeoutException:
        logger.warning("博查搜索超时，回退下一个后端")
        return None
    except Exception as e:
        logger.error(f"博查搜索异常: {e}")
        return None


def _search_with_searxng(query: str, max_results: int) -> Optional[str]:
    """SearXNG 公共实例（免费、无需 API Key）"""
    for instance in _SEARXNG_INSTANCES:
        try:
            client = _get_web_client()
            resp = client.get(
                f"{instance}/search",
                params={
                    "q": query,
                    "format": "json",
                    "language": "zh-CN",
                    "categories": "general",
                    "engines": "google,bing,wikipedia,duckduckgo",
                },
            )
            if resp.status_code != 200:
                continue
            data = resp.json()
            entries = data.get("results", [])
            if not entries:
                continue
            lines = [f"网页搜索 '{query}' 的结果 (via SearXNG):"]
            for i, entry in enumerate(entries[:max_results]):
                title = entry.get("title", "无标题")
                snippet = entry.get("content", entry.get("snippet", ""))
                url = entry.get("url", "")
                lines.append(f"{i + 1}. {title}")
                if snippet:
                    lines.append(f"   {snippet[:200]}")
                if url:
                    lines.append(f"   链接: {url}")
            return "\n".join(lines)
        except httpx.TimeoutException:
            continue
        except Exception:
            continue
    return None


def _search_with_duckduckgo(query: str, max_results: int) -> Optional[str]:
    """DuckDuckGo Instant Answer API（回退）"""
    try:
        client = _get_web_client()
        resp = client.get(
            "https://api.duckduckgo.com/",
            params={"q": query, "format": "json", "no_html": 1, "skip_disambig": 1},
        )
        if resp.status_code != 200:
            return None
        data = resp.json()
        results = []
        if data.get("AbstractText"):
            results.append(f"摘要: {data['AbstractText']}")
        if data.get("RelatedTopics"):
            for topic in data["RelatedTopics"][:max_results]:
                if isinstance(topic, dict) and topic.get("Text"):
                    results.append(f"- {topic['Text']}")
        if results:
            return "\n".join(results)
    except Exception:
        pass
    return None


def _search_web_handler(query: str, max_results: int = 5) -> str:
    """联网搜索 —— 博查（可配置）> SearXNG > DuckDuckGo 三级回退

    优先级：
    1. 博查 API（如果配置了 BOCHA_API_KEY 环境变量）
    2. SearXNG 公共实例（5 个实例轮流尝试）
    3. DuckDuckGo Instant Answer（最终回退）
    """

    # ── 方案 A：博查 Web Search API（需 API Key）──
    result = _search_with_bocha(query, max_results)
    if result:
        return result

    # ── 方案 B：SearXNG 公共实例 ──
    result = _search_with_searxng(query, max_results)
    if result:
        return result

    # ── 方案 C：DuckDuckGo ──
    result = _search_with_duckduckgo(query, max_results)
    if result:
        return result

    return (
        "[联网搜索失败] 所有搜索后端均不可用。\n"
        "建议：1. 检查网络连接  2. 设置 BOCHA_API_KEY 环境变量获取博查 API\n"
        "博查注册: https://open.bochaai.com （免费 2000次/月）"
    )


_search_web_schema = {
    "type": "object",
    "properties": {
        "query": {
            "type": "string",
            "description": "搜索查询词"
        },
        "max_results": {
            "type": "integer",
            "description": "最大返回条数，默认 5",
            "default": 5
        }
    },
    "required": ["query"]
}


def create_search_web_tool() -> ToolDefinition:
    """创建联网搜索工具"""
    return ToolDefinition(
        name="search_web",
        description="联网搜索互联网信息。当你需要获取实时信息、最新资讯、技术文档、行业动态等本地文件之外的知识时使用。使用博查/SearXNG/DuckDuckGo 三级回退，完全免费。",
        parameters=_search_web_schema,
        handler=_search_web_handler,
    )


# ══════════════════════════════════════════════════════════════════════════════
# 预设工具：读取本地文件
# ══════════════════════════════════════════════════════════════════════════════

def _get_allowed_read_roots(db_manager) -> list[Path]:
    """Return canonical paths for explicitly enabled scan directories only."""
    if db_manager is None:
        return []
    from database.models import ScanDirectoryDAO
    return [Path(row['directory_path']).resolve()
            for row in ScanDirectoryDAO(db_manager).get_active()
            if row.get('directory_path')]


def _read_file_handler(file_path: str, max_chars: int = 8000, _db=None) -> str:
    """读取本地文件内容 —— 智能识别类型：.docx/.pdf/.pptx 用专用库提取文本"""
    max_chars = max(1, min(int(max_chars), 8000))
    # The model may only read files from directories the user explicitly scanned.
    try:
        path = Path(file_path).resolve(strict=True)
    except (OSError, RuntimeError) as exc:
        return f"[错误] 文件不存在或无法解析: {exc}"
    if not path.is_file():
        return f"[错误] 不是文件: {file_path}"
    roots = _get_allowed_read_roots(_db)
    if not roots or not any(path.is_relative_to(root) for root in roots):
        return "[安全限制] 只允许读取已启用扫描目录中的文件"

    # Defense in depth for credentials stored inside an otherwise allowed root.
    dangerous_patterns = [
        r'(^|[/\\])\.env$', r'/etc/(passwd|shadow)', r'\\Windows\\System32\\',
        r'\.pem$', r'\.key$', r'id_rsa', r'known_hosts',
    ]
    for pattern in dangerous_patterns:
        if re.search(pattern, file_path, re.IGNORECASE):
            return f"[安全限制] 拒绝读取敏感文件: {file_path}"

    file_path = str(path)
    file_size = path.stat().st_size
    size_str = _format_bytes(file_size)
    if file_size > 10 * 1024 * 1024:  # 超过 10MB 拒绝读取
        return f"[跳过] 文件过大 ({size_str})，拒绝读取内容"

    ext = os.path.splitext(file_path)[1].lower()
    basename = os.path.basename(file_path)

    # ── .docx 文档 ──
    if ext == '.docx':
        try:
            from docx import Document
            doc = Document(file_path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            content = "\n".join(paragraphs)  # 所有段落
            if not content:
                return f"[空文档] {basename} ({size_str}) — 文档中未发现可见文本内容"
            truncated = " ...(内容已截断)" if len(content) > max_chars else ""
            return f"文件: {basename} ({size_str}, 提取 {min(len(content), max_chars)} 字符) [Word 文档]:\n```\n{content[:max_chars]}{truncated}\n```"
        except Exception as e:
            return f"[读取错误] 无法解析 Word 文档 {basename}: {e}"

    # ── .pdf 文档 ──
    if ext == '.pdf':
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(file_path)
            page_count = doc.page_count
            full_text = ""
            for page in doc:
                full_text += page.get_text()
            doc.close()
            if not full_text.strip():
                return f"[空文档] {basename} ({size_str}, {page_count}页) — PDF 中未发现可见文本（可能是扫描版图片 PDF）"
            truncated = " ...(内容已截断)" if len(full_text) > max_chars else ""
            return (
                f"文件: {basename} ({size_str}, {page_count}页, 提取 {min(len(full_text), max_chars)} 字符) [PDF 文档]:\n"
                f"```\n{full_text[:max_chars]}{truncated}\n```"
            )
        except Exception as e:
            return f"[读取错误] 无法解析 PDF 文档 {basename}: {e}"

    # ── .pptx 演示文稿 ──
    if ext == '.pptx':
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            slides = []
            for slide in prs.slides:
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                texts.append(t)
                if texts:
                    slides.append("\n".join(texts))
            content = "\n---\n".join(slides)
            if not content.strip():
                return f"[空文档] {basename} ({size_str}) — PPT 中未发现可见文本"
            truncated = " ...(内容已截断)" if len(content) > max_chars else ""
            return f"文件: {basename} ({size_str}, 提取 {min(len(content), max_chars)} 字符) [PPT 演示文稿]:\n```\n{content[:max_chars]}{truncated}\n```"
        except Exception as e:
            return f"[读取错误] 无法解析 PPT 文件 {basename}: {e}"

    # ── 图片/音频/视频 ──
    binary_exts = {'.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp', '.ico', '.svg',
                   '.mp3', '.wav', '.flac', '.aac', '.ogg', '.wma', '.m4a',
                   '.mp4', '.avi', '.mkv', '.mov', '.wmv', '.webm',
                   '.zip', '.rar', '.7z', '.tar', '.gz',
                   '.exe', '.dll', '.msi', '.apk', '.dmg', '.iso',
                   '.ttf', '.otf', '.woff', '.woff2'}
    if ext in binary_exts:
        return f"[二进制文件] {basename} ({size_str}) — 无法提取文本内容（{ext} 格式）"

    # ── 文本文件（代码、配置、.txt 等） ──
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            content = f.read(max_chars)
    except Exception as e:
        return f"[读取错误] 无法读取文件 {basename}: {e}"

    # 检测是否实际为二进制乱码（虽然读出来了但基本不可打印）
    printable_ratio = sum(1 for c in content if c.isprintable() or c in '\n\r\t ') / max(len(content), 1)
    if printable_ratio < 0.3:
        return f"[二进制文件] {basename} ({size_str}) — 内容为二进制格式，无法提取文本"

    if not content.strip():
        return f"[空文件] {basename} ({size_str})"

    truncated = " ...(内容已截断)" if len(content) >= max_chars else ""
    lang = ext.lstrip(".") if ext else "text"

    return f"文件: {basename} ({size_str}):\n```{lang}\n{content}{truncated}\n```"


_read_file_schema = {
    "type": "object",
    "properties": {
        "file_path": {
            "type": "string",
            "description": "要读取的文件的完整路径"
        },
        "max_chars": {
            "type": "integer",
            "description": "最大返回字符数，默认 8000（约2-3页中文文本）",
            "default": 8000
        }
    },
    "required": ["file_path"]
}


def create_read_file_tool(db_manager=None) -> ToolDefinition:
    """创建文件读取工具"""
    return ToolDefinition(
        name="read_file",
        description="读取本地文件内容并提取文本。支持常见格式：.txt/.md/.py/.js 等文本文件、.docx Word文档、.pdf（含文字）、.pptx 演示文稿。图片/音视频等二进制文件会提示无法读取。",
        parameters=_read_file_schema,
        handler=lambda **kwargs: _read_file_handler(**kwargs, _db=db_manager),
    )


# ══════════════════════════════════════════════════════════════════════════════
# ══════════════════════════════════════════════════════════════════════════════
# 工具函数
# ══════════════════════════════════════════════════════════════════════════════

def _format_bytes(size: int) -> str:
    """格式化字节数为人类可读"""
    if size < 1024:
        return f"{size} B"
    elif size < 1024 * 1024:
        return f"{size / 1024:.1f} KB"
    elif size < 1024 * 1024 * 1024:
        return f"{size / (1024 * 1024):.1f} MB"
    else:
        return f"{size / (1024 * 1024 * 1024):.1f} GB"


def create_default_registry(db_manager=None, ai_layer=None) -> ToolRegistry:
    """创建包含所有预设工具的注册表"""
    registry = ToolRegistry()
    registry.register(create_search_files_tool(db_manager=db_manager, ai_layer=ai_layer))
    registry.register(create_search_content_tool(db_manager=db_manager))
    registry.register(create_search_web_tool())
    registry.register(create_read_file_tool(db_manager=db_manager))
    return registry
