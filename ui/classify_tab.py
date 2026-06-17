"""
分类管理标签页 - 左侧分类树 + 右侧文件列表 + 文件预览面板
"""
from PyQt6.QtWidgets import (
    QWidget, QVBoxLayout, QHBoxLayout, QSplitter, QTreeWidget,
    QTreeWidgetItem, QTableView, QAbstractItemView, QPushButton,
    QLabel, QMessageBox, QHeaderView, QInputDialog, QFileDialog,
    QMenu, QApplication, QProgressBar, QFrame, QTextEdit, QScrollArea,
    QStackedWidget, QSizePolicy, QDialog, QDialogButtonBox, QCheckBox
)
from PyQt6.QtCore import Qt, QThread, pyqtSignal, QSize, QAbstractTableModel, QModelIndex
from PyQt6.QtGui import QColor, QBrush, QPixmap
from PyQt6.QtGui import QAction

import os
import hashlib
from config import FILE_TYPE_NAMES
from core import FileClassifier, FileManager, TagManager
from core.batch_classifier import BatchClassifyWorker
from core.data_cache import GlobalDataCache
from core.rule_engine import TagRecommender
from database.db_manager import db
from database.models import FileDAO, ClassificationDAO, MetadataDAO, TagDAO
from utils.display_utils import format_size, truncate_path, get_file_icon, get_file_color
from utils.logger import logger
from ui.toast import notify
from ui.empty_state import create_empty_state


class BatchOperationWorker(QThread):
    """后台批量操作工作线程（重命名/移动）"""
    progress = pyqtSignal(int, int, str)  # current, total, status
    finished = pyqtSignal(dict)  # results
    error = pyqtSignal(str)

    def __init__(self, operation_func, file_ids: list, extra_args=None, parent=None):
        super().__init__(parent)
        self.operation_func = operation_func
        self.file_ids = file_ids
        self.extra_args = extra_args or {}
        self._cancelled = False

    def cancel(self):
        self._cancelled = True

    def run(self):
        results = {'success': 0, 'failed': 0, 'errors': []}
        total = len(self.file_ids)
        for i, fid in enumerate(self.file_ids):
            if self._cancelled:
                break
            self.progress.emit(i + 1, total, f"处理中 ({i+1}/{total})...")
            try:
                self.operation_func(fid, **self.extra_args)
                results['success'] += 1
            except Exception as e:
                results['failed'] += 1
                results['errors'].append(f"ID {fid}: {e}")
        self.progress.emit(total, total, "完成")
        self.finished.emit(results)


# ── 标签推荐配色方案（与 tags_tab.py 保持一致） ──
_TAG_COLORS = [
    ('#cba6f7', '#1e1e2e'), ('#89b4fa', '#1e1e2e'),
    ('#a6e3a1', '#1e1e2e'), ('#f9e2af', '#1e1e2e'),
    ('#fab387', '#1e1e2e'), ('#94e2d5', '#1e1e2e'),
    ('#f38ba8', '#1e1e2e'), ('#bac2de', '#1e1e2e'),
]
_TAG_LIGHT = [
    ('#cba6f7', '#ffffff'), ('#89b4fa', '#ffffff'),
    ('#a6e3a1', '#1e1e2e'), ('#f9e2af', '#1e1e2e'),
    ('#fab387', '#1e1e2e'), ('#94e2d5', '#1e1e2e'),
    ('#f38ba8', '#ffffff'), ('#bac2de', '#1e1e2e'),
]


def _tag_color_index(name: str) -> int:
    return int(hashlib.md5(name.encode()).hexdigest()[:8], 16) % 8


def _make_tag_btn(text: str, bg: str, fg: str, pt: int = 13, bold: bool = False,
                  checkable: bool = True, checked: bool = False) -> QPushButton:
    """创建与标签云完全一致样式的标签按钮（配色、左对齐、动态圆角、高度）

    Args:
        checked: 按钮初始选中状态，推荐标签默认 False（用户主动勾选）
    """
    btn = QPushButton(text)
    btn.setCheckable(checkable)
    if checkable:
        btn.setChecked(checked)
    btn.setCursor(Qt.CursorShape.PointingHandCursor)
    btn.setFixedHeight(pt + 18)
    btn.setMinimumHeight(32)
    radius = max(8, pt + 2)
    bold_css = 'font-weight: bold;' if bold else ''
    # 显式覆盖全局 QPushButton 的 hover / pressed，防止样式穿透
    hover_qss = (
        f"QPushButton:hover {{ background-color: {bg}; }}"
        f"QPushButton:pressed {{ background-color: {bg}; }}"
    )
    check_qss = (
        f"QPushButton:checked {{ background: {bg}; color: {fg}; border: 2px solid {fg}; }}"
        f"QPushButton:unchecked {{ background: transparent; color: {fg}; border: 1px solid {bg}; }}"
    ) if checkable else ''
    btn.setStyleSheet(
        f"QPushButton {{"
        f"  background-color: {bg};"
        f"  color: {fg};"
        f"  border: none;"
        f"  border-radius: {radius}px;"
        f"  text-align: left;"
        f"  padding: 4px 12px;"
        f"  font-size: {pt}pt;"
        f"  {bold_css}"
        f"}}"
        f"{hover_qss}"
        f"{check_qss}"
    )
    return btn


class DataLoadWorker(QThread):
    """后台数据加载工作线程（在非 UI 线程中进行格式化计算）"""
    data_loaded = pyqtSignal(list, int)  # files, total_count
    load_error = pyqtSignal(str)
    
    def __init__(self, file_dao, mode, current_page, page_size, parent=None):
        super().__init__(parent)
        self.file_dao = file_dao
        self.mode = mode
        self.current_page = current_page
        self.page_size = page_size
    
    def run(self):
        try:
            if self.mode == 'all':
                total_count = self.file_dao.count_active()
                files = self.file_dao.get_all_active_paginated(
                    page=self.current_page, page_size=self.page_size)
            else:
                _, cls_type, cls_value = self.mode
                total_count = self.file_dao.count_by_classification(
                    cls_type, cls_value)
                files = self.file_dao.get_classification_paginated(
                    cls_type, cls_value,
                    page=self.current_page, page_size=self.page_size)
            
            # ── 性能优化：在后台线程预计算所有显示字符串，释放 UI 线程 ──
            for f in files:
                # 文件名（带图标）
                f['_display_name'] = get_file_icon(f['file_type']) + f['file_name']
                # 路径（截断）
                f['_display_path'] = truncate_path(f['file_path'], 60)
                # 类型名
                f['_display_type'] = FILE_TYPE_NAMES.get(f['file_type'], f['file_type'])
                # 类型颜色
                f['_display_color'] = get_file_color(f['file_type'])
                # 文件大小（格式化）
                f['_display_size'] = format_size(f['file_size'])
            
            self.data_loaded.emit(files, total_count)
        except Exception as e:
            self.load_error.emit(str(e))


class FileTableModel(QAbstractTableModel):
    """文件列表数据模型（View-Model 架构，零 QObject 开销）"""
    COLUMNS = ["文件名", "路径", "类型", "大小", "修改时间", "分类"]

    def __init__(self, parent=None):
        super().__init__(parent)
        self._files = []          # 文件记录列表
        self._cls_map = {}        # file_id → [classification_values]
        self._id_to_row = {}      # file_id → row index

    def set_files(self, files: list, cls_map: dict):
        """设置模型数据（调用 beginResetModel/endResetModel 通知 View 全量刷新）"""
        self.beginResetModel()
        self._files = files
        self._cls_map = cls_map
        self._id_to_row = {f['id']: i for i, f in enumerate(files)}
        self.endResetModel()

    def get_record(self, row: int):
        """获取指定行的文件记录"""
        if 0 <= row < len(self._files):
            return self._files[row]
        return None

    def get_record_by_id(self, file_id: int):
        """通过 file_id 获取文件记录"""
        row = self._id_to_row.get(file_id)
        if row is not None:
            return self._files[row]
        return None

    def get_file_id(self, row: int):
        """获取指定行的 file_id"""
        if 0 <= row < len(self._files):
            return self._files[row]['id']
        return None

    def rowCount(self, parent=QModelIndex()):
        return len(self._files) if not parent.isValid() else 0

    def columnCount(self, parent=QModelIndex()):
        return len(self.COLUMNS) if not parent.isValid() else 0

    def headerData(self, section, orientation, role):
        if orientation == Qt.Orientation.Horizontal and role == Qt.ItemDataRole.DisplayRole:
            if 0 <= section < len(self.COLUMNS):
                return self.COLUMNS[section]
        return super().headerData(section, orientation, role)

    def data(self, index, role):
        if not index.isValid():
            return None
        row, col = index.row(), index.column()
        if row < 0 or row >= len(self._files):
            return None
        f = self._files[row]

        if role == Qt.ItemDataRole.UserRole:
            return f['id']

        if role == Qt.ItemDataRole.DisplayRole:
            if col == 0:
                return f.get('_display_name')
            elif col == 1:
                return f.get('_display_path')
            elif col == 2:
                return f.get('_display_type')
            elif col == 3:
                return f.get('_display_size')
            elif col == 4:
                mtime = f.get('modify_time', '')
                return str(mtime) if mtime else ""
            elif col == 5:
                cls_values = self._cls_map.get(f['id'], [])
                return ", ".join(cls_values) if cls_values else "-"

        if role == Qt.ItemDataRole.ToolTipRole:
            if col == 0:
                return f['file_name']
            elif col == 1:
                return f['file_path']

        if role == Qt.ItemDataRole.ForegroundRole and col == 2:
            color = f.get('_display_color')
            if color:
                return QColor(color)

        return None


class TreeLoadWorker(QThread):
    """后台分类树加载工作线程"""
    tree_loaded = pyqtSignal(dict)  # tree_data
    
    def __init__(self, classifier, parent=None):
        super().__init__(parent)
        self.classifier = classifier
    
    def run(self):
        try:
            tree_data = self.classifier.get_classification_tree()
            self.tree_loaded.emit(tree_data)
        except Exception:
            self.tree_loaded.emit({})


class ClassifyTab(QWidget):
    def __init__(self, parent=None):
        super().__init__(parent)
        self.file_dao = FileDAO(db)
        self.cls_dao = ClassificationDAO(db)
        self.meta_dao = MetadataDAO(db)
        self.tag_manager = TagManager()
        self.classifier = FileClassifier()
        self.file_manager = FileManager()
        self.page_size = 500
        self.current_page = 0
        self._total_count = 0
        # 当前加载模式：'all' 或 ('classification', cls_type, cls_value)
        self._mode = 'all'
        # 当前预览文件的路径
        self._current_preview_path = None
        # 缓存分类树数据，避免每次刷新都重新查询
        self._tree_cache = None
        self._tree_cache_time = 0
        # 是否正在加载数据
        self._is_loading = False
        # 后台数据加载工作线程
        self._data_worker = None
        # 后台分类树加载工作线程
        self._tree_worker = None
        # 文件列表数据缓存：{mode_key: (files, total_count, timestamp)}
        self._data_cache = {}
        self._data_cache_max_age = 30  # 缓存有效期（秒）
        # 全局数据缓存服务
        self._global_cache = GlobalDataCache.get_instance()
        self._init_ui()

    # ──── 多文件类型预览：扩展名与工具方法 ────

    CODE_EXTS = {
        '.py', '.js', '.ts', '.jsx', '.tsx', '.java', '.cpp', '.c', '.h', '.hpp',
        '.html', '.htm', '.css', '.scss', '.less', '.sql', '.php', '.rb', '.go',
        '.rs', '.swift', '.kt', '.sh', '.bat', '.ps1', '.vue', '.svelte',
        '.xml', '.yaml', '.yml', '.toml', '.json', '.ini', '.cfg', '.r', '.lua',
        '.dart', '.scala', '.pl', '.m', '.mm',
    }

    TEXT_EXTS = {'.txt', '.csv', '.log', '.md', '.rst', '.readme', '.tsv'}

    PDF_EXTS = {'.pdf'}

    DOCX_EXTS = {'.docx', '.pptx'}

    DOC_EXTS = {'.doc', '.ppt'}

    @staticmethod
    def _is_text_file(ext):
        ext = ext.lower()
        return ext in ClassifyTab.CODE_EXTS or ext in ClassifyTab.TEXT_EXTS

    @staticmethod
    def _detect_encoding(file_path):
        encodings = ['utf-8-sig', 'utf-8', 'gbk', 'gb18030', 'gb2312', 'latin-1']
        for enc in encodings:
            try:
                with open(file_path, 'r', encoding=enc) as f:
                    f.read(4096)
                return enc
            except (UnicodeDecodeError, UnicodeError):
                continue
        return 'utf-8'

    @staticmethod
    def _read_file_content(file_path, max_chars=100000):
        encoding = ClassifyTab._detect_encoding(file_path)
        try:
            with open(file_path, 'r', encoding=encoding, errors='replace') as f:
                content = f.read(max_chars)
            return content, encoding
        except Exception:
            return None, None

    def _init_ui(self):
        layout = QVBoxLayout(self)
        layout.setContentsMargins(16, 16, 16, 16)
        layout.setSpacing(8)

        # 顶部操作栏
        top_layout = QHBoxLayout()

        reclassify_btn = QPushButton("重新分类所有文件")
        reclassify_btn.setObjectName("primaryBtn")
        reclassify_btn.clicked.connect(self._reclassify_all)
        top_layout.addWidget(reclassify_btn)

        top_layout.addStretch()

        self.file_count_label = QLabel("")
        self.file_count_label.setObjectName("subtitleLabel")
        top_layout.addWidget(self.file_count_label)

        layout.addLayout(top_layout)

        # 分割器: 左树 + 中列表 + 右预览
        splitter = QSplitter(Qt.Orientation.Horizontal)

        # 左侧分类树
        left_widget = QWidget()
        left_layout = QVBoxLayout(left_widget)
        left_layout.setContentsMargins(0, 0, 0, 0)

        tree_title = QLabel("分类导航")
        tree_title.setStyleSheet("font-weight: bold; color: #cba6f7; padding: 4px;")
        left_layout.addWidget(tree_title)

        self.tree = QTreeWidget()
        self.tree.setHeaderHidden(True)
        self.tree.itemClicked.connect(self._on_tree_click)
        left_layout.addWidget(self.tree)

        splitter.addWidget(left_widget)

        # 中间文件列表
        center_widget = QWidget()
        right_layout = QVBoxLayout(center_widget)
        right_layout.setContentsMargins(0, 0, 0, 0)

        # 批量操作
        batch_layout = QHBoxLayout()

        rename_btn = QPushButton("批量重命名")
        rename_btn.clicked.connect(self._batch_rename)
        batch_layout.addWidget(rename_btn)

        move_btn = QPushButton("批量移动")
        move_btn.clicked.connect(self._batch_move)
        batch_layout.addWidget(move_btn)

        batch_layout.addStretch()

        self.selected_label = QLabel("")
        self.selected_label.setObjectName("subtitleLabel")
        batch_layout.addWidget(self.selected_label)

        right_layout.addLayout(batch_layout)

        self.file_table = QTableView()
        self.file_table.setSelectionBehavior(QAbstractItemView.SelectionBehavior.SelectRows)
        self.file_table.setSelectionMode(QAbstractItemView.SelectionMode.MultiSelection)
        self.file_table.setAlternatingRowColors(True)
        self.file_table.setSortingEnabled(True)
        self.file_table.setContextMenuPolicy(Qt.ContextMenuPolicy.CustomContextMenu)
        self.file_table.customContextMenuRequested.connect(self._show_context_menu)

        # 创建数据模型并绑定
        self._file_model = FileTableModel(self)
        self.file_table.setModel(self._file_model)

        # 列宽设置
        header = self.file_table.horizontalHeader()
        header.setSectionResizeMode(0, QHeaderView.ResizeMode.Interactive)
        header.setSectionResizeMode(1, QHeaderView.ResizeMode.Stretch)
        header.setSectionResizeMode(2, QHeaderView.ResizeMode.ResizeToContents)
        header.setSectionResizeMode(3, QHeaderView.ResizeMode.ResizeToContents)
        header.setSectionResizeMode(4, QHeaderView.ResizeMode.ResizeToContents)
        header.setSectionResizeMode(5, QHeaderView.ResizeMode.ResizeToContents)
        self.file_table.setColumnWidth(0, 180)

        # 选中信号改用 selectionModel
        self.file_table.selectionModel().selectionChanged.connect(
            self._on_selection_changed)
        right_layout.addWidget(self.file_table)

        # 空状态引导
        self._empty_state = create_empty_state('classify', parent=center_widget)
        right_layout.addWidget(self._empty_state)

        # 分页控件
        page_layout = QHBoxLayout()
        self.prev_page_btn = QPushButton("上一页")
        self.prev_page_btn.clicked.connect(self._prev_page)
        self.prev_page_btn.setEnabled(False)
        page_layout.addWidget(self.prev_page_btn)

        page_layout.addStretch()
        self.page_label = QLabel("第 1 页 / 共 1 页")
        page_layout.addWidget(self.page_label)
        page_layout.addStretch()

        self.next_page_btn = QPushButton("下一页")
        self.next_page_btn.clicked.connect(self._next_page)
        self.next_page_btn.setEnabled(False)
        page_layout.addWidget(self.next_page_btn)

        right_layout.addLayout(page_layout)

        # 分类进度条（默认隐藏）
        self.reclassify_progress = QProgressBar()
        self.reclassify_progress.setVisible(False)
        right_layout.addWidget(self.reclassify_progress)

        self.reclassify_label = QLabel("")
        self.reclassify_label.setObjectName("subtitleLabel")
        self.reclassify_label.setVisible(False)
        right_layout.addWidget(self.reclassify_label)

        splitter.addWidget(center_widget)

        # 右侧预览面板（默认折叠隐藏）
        self._preview_panel = self._build_preview_panel()
        self._preview_panel.setVisible(False)
        splitter.addWidget(self._preview_panel)
        self._splitter = splitter

        splitter.setSizes([180, 820, 0])

        layout.addWidget(splitter, 1)

    def refresh_data(self):
        """刷新数据：重建分类树 + 重新加载当前页"""
        # 避免重复加载
        if self._is_loading:
            return
        self._is_loading = True
        
        # 清除全局缓存（数据可能已变化）
        self._global_cache.invalidate_cache()
        
        # 显示加载状态
        self.file_count_label.setText("加载中...")
        
        # 清除树缓存，强制重新加载
        self._tree_cache = None
        self._tree_cache_time = 0
        
        # 异步加载分类树
        self._build_tree_async()
        
        # 异步加载文件列表
        self._reload_page()
    
    def _build_tree_async(self):
        """异步构建分类树"""
        import time
        current_time = time.time()
        
        # 检查缓存是否有效（30秒内）
        if self._tree_cache is not None and (current_time - self._tree_cache_time) < 30:
            # 使用缓存数据，直接构建树（很快）
            self._build_tree_from_data(self._tree_cache)
            return
        
        # 停止之前的 tree worker
        if hasattr(self, '_tree_worker') and self._tree_worker is not None:
            self._tree_worker.quit()
            self._tree_worker.wait()
        
        # 创建后台工作线程加载树
        self._tree_worker = TreeLoadWorker(self.classifier, self)
        self._tree_worker.tree_loaded.connect(self._on_tree_loaded)
        self._tree_worker.start()
    
    def _on_tree_loaded(self, tree_data: dict):
        """分类树加载完成的回调"""
        import time
        self._tree_cache = tree_data
        self._tree_cache_time = time.time()
        self._build_tree_from_data(tree_data)
        self._tree_worker = None
    
    def _build_tree_from_data(self, tree_data: dict):
        """根据数据构建分类树（在 UI 线程中执行，很快）"""
        self.tree.clear()

        # 全部文件
        all_item = QTreeWidgetItem(self.tree, ["全部文件"])
        all_item.setData(0, Qt.ItemDataRole.UserRole, ('all', None))

        # 按维度分组：按类型、按日期、按关键词
        dimension_names = {'by_type': '按类型', 'by_date': '按日期', 'by_keyword': '按关键词'}
        
        for category, items in tree_data.items():
            # 创建维度父节点（如"按类型"、"按日期"）
            parent = QTreeWidgetItem(self.tree, [category])
            parent.setExpanded(True)
            
            # 去重：同一个分类值只显示一次
            seen_values = set()
            for value, count in items:
                if value in seen_values:
                    continue
                seen_values.add(value)
                
                child = QTreeWidgetItem(parent, [f"{value} ({count})"])
                # 存储分类类型和值
                db_type = {'按类型': 'by_type', '按日期': 'by_date', '按关键词': 'by_keyword'}.get(category, category)
                child.setData(0, Qt.ItemDataRole.UserRole, (db_type, value))

        self.tree.expandAll()

    def _on_tree_click(self, item, column):
        """分类树点击事件"""
        data = item.data(0, Qt.ItemDataRole.UserRole)
        if not data:
            return
        
        # 避免重复点击相同分类
        category, value = data
        if category == 'all':
            new_mode = 'all'
        else:
            # category 已经是 db_type（by_type, by_date, by_keyword）
            new_mode = ('classification', category, value)
        
        # 如果点击的是当前已选中的分类，不重复加载
        if new_mode == self._mode and self.current_page == 0:
            return
        
        self.current_page = 0
        self._mode = new_mode
        self._reload_page()

    def _get_mode_key(self) -> str:
        """获取当前模式的缓存键"""
        if self._mode == 'all':
            return f"all_page_{self.current_page}"
        else:
            _, cls_type, cls_value = self._mode
            return f"{cls_type}_{cls_value}_page_{self.current_page}"
    
    def _get_cached_data(self) -> tuple:
        """尝试从缓存获取数据，返回 (files, total_count] 或 None"""
        import time
        mode_key = self._get_mode_key()
        if mode_key in self._data_cache:
            files, total_count, timestamp = self._data_cache[mode_key]
            age = time.time() - timestamp
            if age < self._data_cache_max_age:
                return (files, total_count)
        return None
    
    def _cache_data(self, files: list, total_count: int):
        """缓存当前数据"""
        import time
        mode_key = self._get_mode_key()
        self._data_cache[mode_key] = (files, total_count, time.time())
        
        # 限制缓存大小，保留最近 20 个
        if len(self._data_cache) > 20:
            # 删除最旧的缓存
            oldest_key = min(self._data_cache.keys(), 
                           key=lambda k: self._data_cache[k][2])
            del self._data_cache[oldest_key]
    
    def _reload_page(self):
        """根据当前模式从数据库加载当前页（优先全局缓存）"""
        # 1. 优先从全局缓存读取（预加载的数据，0 延迟）
        global_cached = self._global_cache.get_files(
            self._mode, self.current_page, self.page_size)
        if global_cached is not None:
            files, total_count = global_cached
            self._total_count = total_count
            self._populate_table(files)
            # 缓存命中，无需后台查询
            return
        
        # 2. 全局缓存未命中，尝试本地缓存
        cached = self._get_cached_data()
        if cached is not None:
            files, total_count = cached
            self._total_count = total_count
            self._populate_table(files)
            # 显示"数据来自缓存"提示
            self.file_count_label.setText(f"共 {total_count} 个文件（正在刷新...）")
        else:
            # 无缓存，显示加载状态
            self.file_count_label.setText("加载中...")
            self.file_table.setVisible(False)
            self._empty_state.setVisible(False)
        
        # 3. 停止之前的 worker（如果存在）
        if hasattr(self, '_data_worker') and self._data_worker is not None:
            self._data_worker.quit()
            self._data_worker.wait()
        
        # 4. 创建后台工作线程加载最新数据
        self._data_worker = DataLoadWorker(
            self.file_dao, self._mode, self.current_page, self.page_size, self)
        self._data_worker.data_loaded.connect(self._on_data_loaded)
        self._data_worker.load_error.connect(self._on_load_error)
        self._data_worker.start()
    
    def _on_data_loaded(self, files: list, total_count: int):
        """数据加载完成的回调"""
        # 缓存数据到本地和全局缓存
        self._cache_data(files, total_count)
        self._global_cache.cache_files(
            self._mode, self.current_page, files, total_count)
        
        # 更新显示
        self._total_count = total_count
        self._populate_table(files)
        self._data_worker = None
        self._is_loading = False
    
    def _on_load_error(self, error_msg: str):
        """数据加载失败的回调"""
        logger.error(f"加载文件失败: {error_msg}")
        self.file_count_label.setText("加载失败")
        self.file_table.setVisible(False)
        self._empty_state.setVisible(True)
        self._data_worker = None
        self._is_loading = False

    def _populate_table(self, files):
        # 修复：翻页/切换分类时清除上一页的选中状态，避免跨页选中残留
        self.file_table.clearSelection()
        total = self._total_count
        total_pages = max(1, (total + self.page_size - 1) // self.page_size)

        # 修正当前页范围
        if self.current_page >= total_pages:
            self.current_page = total_pages - 1

        self.file_count_label.setText(f"共 {total} 个文件")

        # 空状态检测
        self._empty_state.setVisible(len(files) == 0)
        self.file_table.setVisible(len(files) > 0)

        # 批量查询分类（避免 N+1）
        page_ids = [f['id'] for f in files]
        try:
            cls_map = self.cls_dao.get_by_file_ids(page_ids)
        except Exception:
            cls_map = {}

        # ── 直接设置模型数据，零 Q表单项创建 ──
        self._file_model.set_files(files, cls_map)

        # 更新分页状态
        self.page_label.setText(f"第 {self.current_page + 1} 页 / 共 {total_pages} 页")
        self.prev_page_btn.setEnabled(self.current_page > 0)
        self.next_page_btn.setEnabled(self.current_page < total_pages - 1)

    def _prev_page(self):
        if self.current_page > 0:
            self.current_page -= 1
            self._reload_page()

    def _next_page(self):
        total_pages = max(1, (self._total_count + self.page_size - 1) // self.page_size)
        if self.current_page < total_pages - 1:
            self.current_page += 1
            self._reload_page()

    def _on_selection_changed(self, selected=None, deselected=None):
        count = len(self.file_table.selectionModel().selectedRows())
        self.selected_label.setText(f"已选 {count} 个文件" if count > 0 else "")
        # 更新预览面板
        if count == 1:
            rows = self.file_table.selectionModel().selectedRows()
            if rows:
                fid = self._file_model.data(rows[0], Qt.ItemDataRole.UserRole)
                if fid is not None:
                    self._show_preview(fid)
                    return
        self._clear_preview()

    def _show_context_menu(self, pos):
        """文件列表右键菜单（通过 Model 获取数据，零数据库查询）"""
        row = self.file_table.rowAt(pos.y())
        if row < 0:
            return
        file_id = self._file_model.get_file_id(row)
        if file_id is None:
            return

        # 从模型缓存获取 record
        record = self._file_model.get_record(row)
        if not record:
            return

        file_path = record.get('file_path', '')
        file_path = os.path.normpath(file_path) if isinstance(file_path, str) and file_path else ''

        logger.debug(f"右键菜单 - 文件路径: {file_path}")

        menu = QMenu(self)

        # 修复：用具名函数 + 显式参数，避免 PyQt6 QAction.triggered(bool) 默认参数陷阱
        def _do_open_file(checked=False, fp=file_path, fid=file_id):
            self._safe_open_file(fp, file_id=fid)

        def _do_open_folder(checked=False, fp=file_path, fid=file_id):
            self._safe_open_folder(fp, file_id=fid)

        open_action = QAction("打开文件", self)
        open_action.triggered.connect(_do_open_file)
        menu.addAction(open_action)

        open_folder_action = QAction("打开所在文件夹", self)
        open_folder_action.triggered.connect(_do_open_folder)
        menu.addAction(open_folder_action)

        menu.addSeparator()

        copy_path_action = QAction("复制路径", self)
        copy_path_action.triggered.connect(
            lambda: QApplication.clipboard().setText(file_path))
        menu.addAction(copy_path_action)

        copy_name_action = QAction("复制文件名", self)
        copy_name_action.triggered.connect(
            lambda: QApplication.clipboard().setText(record['file_name']))
        menu.addAction(copy_name_action)

        menu.addSeparator()

        rename_action = QAction("重命名", self)
        rename_action.triggered.connect(lambda: self._context_rename(file_id))
        menu.addAction(rename_action)

        # 智能推荐标签
        recommend_tag_action = QAction("智能推荐标签", self)
        recommend_tag_action.triggered.connect(
            lambda: self._recommend_tags_for_file(file_id))
        menu.addAction(recommend_tag_action)

        menu.addSeparator()

        delete_action = QAction("标记删除", self)
        delete_action.triggered.connect(lambda: self._context_delete(file_id))
        menu.addAction(delete_action)

        permanent_delete_action = QAction("永久删除", self)
        permanent_delete_action.setObjectName("dangerBtn")
        permanent_delete_action.triggered.connect(
            lambda: self._context_permanent_delete(file_id))
        menu.addAction(permanent_delete_action)

        menu.exec(self.file_table.viewport().mapToGlobal(pos))

    def _safe_open_file(self, file_path, file_id=None):
        """安全打开文件，文件不存在时给出友好提示

        修复：兼容 file_path 为空/False/非字符串等异常情况，
        优先尝试用 file_id 从数据库反查真实路径。
        """
        # 路径无效判定（兼容空串、False 字符串、None 等异常值）
        BAD_VALUES = ('', 'False', 'True', 'false', 'true', '0', 'null', 'None')
        needs_lookup = (
            file_path is None
            or not isinstance(file_path, str)
            or file_path.strip() in BAD_VALUES
        )

        if needs_lookup and file_id is not None:
            try:
                rec = self.file_dao.get_by_id(file_id)
                if rec and rec.get('file_path'):
                    file_path = rec['file_path']
                    logger.info(f"路径无效已用 file_id={file_id} 反查: {file_path!r}")
            except Exception as e:
                logger.warning(f"反查 file_id={file_id} 失败: {e}")

        # 再次校验
        if (file_path is None
                or not isinstance(file_path, str)
                or file_path.strip() in BAD_VALUES):
            notify(self, "无法操作：文件路径无效", 'warning', 4000)
            logger.warning(
                f"文件路径无效: {file_path!r} (type: {type(file_path).__name__})")
            return

        logger.debug(f"尝试打开文件: {file_path}")
        file_path = os.path.normpath(file_path)
        logger.debug(f"标准化后路径: {file_path}")

        if not os.path.exists(file_path):
            notify(self, f"文件不存在: {os.path.basename(file_path)}", 'warning', 4000)
            logger.warning(f"尝试打开不存在的文件: {file_path}")
            return

        try:
            os.startfile(file_path)
            logger.info(f"成功打开文件: {file_path}")
        except Exception as e:
            notify(self, f"无法打开文件: {e}", 'error', 5000)
            logger.error(f"打开文件失败: {file_path}, 错误: {e}")

    def _safe_open_folder(self, file_path, file_id=None):
        """安全打开所在文件夹

        修复：兼容 file_path 为空/False/非字符串等异常情况，
        优先尝试用 file_id 从数据库反查真实路径。
        """
        BAD_VALUES = ('', 'False', 'True', 'false', 'true', '0', 'null', 'None')
        needs_lookup = (
            file_path is None
            or not isinstance(file_path, str)
            or file_path.strip() in BAD_VALUES
        )

        if needs_lookup and file_id is not None:
            try:
                rec = self.file_dao.get_by_id(file_id)
                if rec and rec.get('file_path'):
                    file_path = rec['file_path']
                    logger.info(f"路径无效已用 file_id={file_id} 反查: {file_path!r}")
            except Exception as e:
                logger.warning(f"反查 file_id={file_id} 失败: {e}")

        if (file_path is None
                or not isinstance(file_path, str)
                or file_path.strip() in BAD_VALUES):
            notify(self, "无法操作：文件路径无效", 'warning', 4000)
            logger.warning(
                f"文件路径无效: {file_path!r} (type: {type(file_path).__name__})")
            return

        logger.debug(f"尝试打开文件夹: {file_path}")
        file_path = os.path.normpath(file_path)
        folder = os.path.dirname(file_path)
        logger.debug(f"文件夹路径: {folder}")

        if not folder or not os.path.exists(folder):
            notify(self, "文件所在目录已不存在", 'warning', 4000)
            logger.warning(f"文件夹不存在: {folder}")
            return

        try:
            os.startfile(folder)
            logger.info(f"成功打开文件夹: {folder}")
        except Exception as e:
            notify(self, f"无法打开文件夹: {e}", 'error', 5000)
            logger.error(f"打开文件夹失败: {folder}, 错误: {e}")

    def _context_rename(self, file_id):
        """右键菜单：重命名单个文件（需二次确认）"""
        record = self._file_model.get_record_by_id(file_id)
        if record is None:
            record = self.file_dao.get_by_id(file_id)
        if not record:
            return
        new_name, ok = QInputDialog.getText(
            self, "重命名", "输入新文件名:",
            text=record['file_name'])
        if not ok or not new_name:
            return
        reply = QMessageBox.question(
            self, "确认重命名",
            f"确定要将\n{record['file_name']}\n重命名为\n{new_name}?\n\n"
            "此操作会真的修改硬盘上的文件名！")
        if reply == QMessageBox.StandardButton.Yes:
            try:
                self.file_manager.rename_file(file_id, new_name=new_name)
                self.refresh_data()
                notify(self, f"已重命名为: {new_name}", 'success', 3000)
            except Exception as e:
                QMessageBox.critical(self, "重命名失败", str(e))

    def _recommend_tags_for_file(self, file_id):
        """智能推荐标签：弹出标签云风格对话框，点击彩色标签即可选中/取消"""
        record = self._file_model.get_record_by_id(file_id)
        if record is None:
            record = self.file_dao.get_by_id(file_id)
        if not record:
            notify(self, "文件记录不存在", 'warning', 3000)
            return

        # 获取推荐标签
        recommendations = TagRecommender.recommend(record, max_tags=8)
        if not recommendations:
            QMessageBox.information(self, "标签推荐",
                                    f"未找到适合「{record['file_name']}」的推荐标签。")
            return

        # 获取已有标签，避免重复推荐
        existing_tags = set(t['tag_name'] for t in self.tag_manager.get_tags_by_file(file_id))

        # 构建推荐标签列表（排除已有标签）
        suggested = [(tag, conf) for tag, conf in recommendations if tag not in existing_tags]
        if not suggested:
            QMessageBox.information(self, "标签推荐",
                                    f"「{record['file_name']}」已有标签覆盖了所有推荐。")
            return

        # ── 构建标签云风格对话框 ──
        is_light = getattr(self, '_theme', 'dark') == 'light'
        dialog_bg = '#eff1f5' if is_light else '#1e1e2e'
        text_fg = '#4c4f69' if is_light else '#cdd6f4'
        subtitle_fg = '#7c7f93' if is_light else '#a6adc8'
        sep_color = '#ccd0da' if is_light else '#45475a'

        dlg = QDialog(self)
        dlg.setWindowTitle("智能推荐标签")
        dlg.setMinimumWidth(420)
        dlg.setModal(True)
        dlg.setStyleSheet(f"QDialog {{ background-color: {dialog_bg}; }}")

        layout = QVBoxLayout(dlg)
        layout.setSpacing(12)

        # 文件信息
        info = QLabel(f"为「{record['file_name']}」推荐以下标签（点击切换选中/取消）：")
        info.setWordWrap(True)
        info.setStyleSheet(f"font-size: 12px; color: {text_fg}; margin-bottom: 4px;")
        layout.addWidget(info)

        # 已有标签展示（使用与标签云一致的按钮样式，不可切换）
        if existing_tags:
            existing_palette = _TAG_LIGHT if is_light else _TAG_COLORS
            existing_row = QHBoxLayout()
            existing_row.setSpacing(6)
            el = QLabel("已有:")
            el.setStyleSheet(f"color: {subtitle_fg}; font-size: 11pt;")
            existing_row.addWidget(el)
            for et in sorted(existing_tags):
                idx = _tag_color_index(et)
                bg, fg = existing_palette[idx]
                tag_btn = _make_tag_btn(f"  {et}  ", bg, fg, pt=13, bold=False, checkable=False)
                existing_row.addWidget(tag_btn)
            existing_row.addStretch()
            layout.addLayout(existing_row)

        # 分隔线
        sep = QFrame()
        sep.setFrameShape(QFrame.Shape.HLine)
        sep.setStyleSheet(f"color: {sep_color};")
        layout.addWidget(sep)

        # 推荐标签按钮（标签云样式，与 tags_tab.py 配色一致）
        palette = _TAG_LIGHT if is_light else _TAG_COLORS
        tag_buttons = []

        tags_layout = QVBoxLayout()
        tags_layout.setSpacing(8)

        for tag, conf in suggested:
            idx = _tag_color_index(tag)
            bg, fg = palette[idx]
            display = f"{tag}  ({conf:.0%})"
            btn = _make_tag_btn(display, bg, fg, pt=13, bold=False, checkable=True)
            tag_buttons.append((tag, btn, conf))
            tags_layout.addWidget(btn)

        layout.addLayout(tags_layout)

        # 按钮区
        btn_box = QDialogButtonBox(QDialogButtonBox.StandardButton.Ok | QDialogButtonBox.StandardButton.Cancel)
        btn_box.accepted.connect(dlg.accept)
        btn_box.rejected.connect(dlg.reject)
        layout.addWidget(btn_box)

        # 显示对话框
        if dlg.exec() != QDialog.DialogCode.Accepted:
            return

        # 收集选中的标签
        selected = [tag for tag, btn, _ in tag_buttons if btn.isChecked()]
        if not selected:
            return

        try:
            self.tag_manager.batch_add_tags([file_id], selected)
            notify(self, f"已添加标签: {', '.join(selected)}", 'success', 3500)
            logger.info(f"标签推荐: file_id={file_id} 添加了标签 {selected}")
        except Exception as e:
            logger.error(f"添加标签失败: {e}")
            QMessageBox.critical(self, "标签推荐", f"添加标签失败: {e}")

    def _context_delete(self, file_id):
        """右键菜单：标记删除单个文件"""
        reply = QMessageBox.question(
            self, "确认删除", "确定标记删除该文件?")
        if reply == QMessageBox.StandardButton.Yes:
            try:
                self.file_manager.delete_file(file_id)
                self.refresh_data()
                notify(self, "文件已标记删除", 'success', 3000)
            except Exception as e:
                QMessageBox.critical(self, "删除失败", str(e))

    def _context_permanent_delete(self, file_id):
        """右键菜单：永久删除文件（从硬盘清除）"""
        record = self._file_model.get_record_by_id(file_id)
        if record is None:
            record = self.file_dao.get_by_id(file_id)
        if not record:
            return
        reply = QMessageBox.question(
            self, "⚠️ 永久删除",
            f"确定要永久删除以下文件?\n\n{record['file_name']}\n\n"
            "此操作将从硬盘上彻底删除文件，不可恢复！\n"
            "（操作历史中的撤销不可用于永久删除）",
            QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No,
            QMessageBox.StandardButton.No)
        if reply == QMessageBox.StandardButton.Yes:
            # 二次确认
            reply2 = QMessageBox.question(
                self, "⚠️ 最终确认",
                "此操作无法撤销，确定继续？",
                QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No,
                QMessageBox.StandardButton.No)
            if reply2 == QMessageBox.StandardButton.Yes:
                try:
                    self.file_manager.permanent_delete(file_id)
                    self.refresh_data()
                except Exception as e:
                    QMessageBox.critical(self, "删除失败", str(e))

    def rename_selected(self):
        """F2 重命名当前选中的第一个文件"""
        rows = self.file_table.selectionModel().selectedRows()
        if not rows:
            notify(self, "请先选择要重命名的文件", 'info', 2000)
            return
        fid = self._file_model.data(rows[0], Qt.ItemDataRole.UserRole)
        if fid is not None:
            self._context_rename(fid)

    def delete_selected(self):
        """Delete 标记删除选中的文件"""
        ids = self._get_selected_ids()
        if not ids:
            notify(self, "请先选择要删除的文件", 'info', 2000)
            return
        reply = QMessageBox.question(
            self, "确认删除", f"确定标记删除选中的 {len(ids)} 个文件?")
        if reply == QMessageBox.StandardButton.Yes:
            success = 0
            for fid in ids:
                try:
                    self.file_manager.delete_file(fid)
                    success += 1
                except Exception as e:
                    logger.warning(f"删除失败 ID={fid}: {e}")
            notify(self, f"已标记删除 {success} 个文件", 'success', 3000)
            self.refresh_data()

    def _get_selected_ids(self):
        """从表格选中行读取文件ID（通过 Model，不受排序影响）"""
        ids = []
        for idx in self.file_table.selectionModel().selectedRows():
            fid = self._file_model.data(idx, Qt.ItemDataRole.UserRole)
            if fid is not None:
                ids.append(fid)
        return ids

    def _batch_rename(self):
        ids = self._get_selected_ids()
        if not ids:
            QMessageBox.information(self, "提示", "请先选择要重命名的文件")
            return
        reply = QMessageBox.question(
            self, "确认批量重命名",
            f"将对 {len(ids)} 个文件执行重命名\n格式: 日期_类型_原名\n确定继续?")
        if reply != QMessageBox.StandardButton.Yes:
            return
        self._start_batch_operation(
            operation_func=self.file_manager.rename_file,
            file_ids=ids,
            operation_name="批量重命名")

    def _batch_move(self):
        ids = self._get_selected_ids()
        if not ids:
            QMessageBox.information(self, "提示", "请先选择要移动的文件")
            return
        target = QFileDialog.getExistingDirectory(self, "选择目标目录")
        if not target:
            return
        self._start_batch_operation(
            operation_func=self.file_manager.move_file,
            file_ids=ids,
            operation_name="批量移动",
            extra_args={'target_dir': target})

    def _start_batch_operation(self, operation_func, file_ids, operation_name, extra_args=None):
        """启动后台批量操作，显示进度"""
        self.reclassify_progress.setVisible(True)
        self.reclassify_label.setVisible(True)
        self.reclassify_progress.setValue(0)
        self.reclassify_progress.setMaximum(len(file_ids))
        self.reclassify_label.setText(f"{operation_name} 准备中...")

        self._batch_worker = BatchOperationWorker(
            operation_func=operation_func,
            file_ids=file_ids,
            extra_args=extra_args)
        self._batch_worker.progress.connect(self._on_batch_progress)
        self._batch_worker.finished.connect(
            lambda results, name=operation_name: self._on_batch_finished(results, name))
        self._batch_worker.start()

    def _on_batch_progress(self, current, total, status):
        self.reclassify_progress.setMaximum(total)
        self.reclassify_progress.setValue(current)
        self.reclassify_label.setText(status)

    def _on_batch_finished(self, results, operation_name):
        self.reclassify_progress.setVisible(False)
        self.reclassify_label.setVisible(False)
        msg = f"{operation_name}完成: 成功 {results['success']}, 失败 {results['failed']}"
        notify(self, msg, 'success' if results['failed'] == 0 else 'warning', 4000)
        self.refresh_data()

    def _reclassify_all(self):
        reply = QMessageBox.question(self, "确认", "重新分类所有文件? 这将清除现有分类结果。")
        if reply != QMessageBox.StandardButton.Yes:
            return

        # 禁用按钮，显示进度
        self.reclassify_progress.setVisible(True)
        self.reclassify_label.setVisible(True)
        self.reclassify_progress.setValue(0)
        self.reclassify_label.setText("正在重新分类...")

        self.worker = BatchClassifyWorker(self.classifier)
        self.worker.progress.connect(self._on_reclassify_progress)
        self.worker.finished.connect(self._on_reclassify_finished)
        self.worker.error.connect(self._on_reclassify_error)
        self.worker.start()

    def _on_reclassify_progress(self, current, total):
        self.reclassify_progress.setMaximum(total)
        self.reclassify_progress.setValue(current)
        self.reclassify_label.setText(f"正在分类 ({current}/{total})...")

    def _on_reclassify_finished(self, count):
        self.reclassify_progress.setVisible(False)
        self.reclassify_label.setVisible(False)
        notify(self, f"分类完成: 已分类 {count} 个文件", 'success', 4000)
        self.refresh_data()

    def _on_reclassify_error(self, msg):
        self.reclassify_progress.setVisible(False)
        self.reclassify_label.setVisible(False)
        notify(self, f"分类失败: {msg}", 'error', 5000)

    # ── 文件预览面板 ──

    def _build_preview_panel(self) -> QWidget:
        """构建右侧文件预览面板（支持图片/文本/代码/PDF）"""
        panel = QFrame()
        panel.setMinimumWidth(230)
        panel.setFrameShape(QFrame.Shape.StyledPanel)
        panel.setStyleSheet(
            "QFrame { border: 1px solid #313244; border-radius: 6px; }")

        self._preview_layout = QVBoxLayout(panel)
        self._preview_layout.setContentsMargins(10, 10, 10, 10)
        self._preview_layout.setSpacing(6)

        # 标题
        self._preview_title = QLabel("文件预览")
        self._preview_title.setStyleSheet(
            "font-weight: bold; font-size: 12px;"
            "border: none; background: transparent;")
        self._preview_layout.addWidget(self._preview_title)

        # 多类型预览栈
        self._preview_stack = QStackedWidget()
        self._preview_stack.setSizePolicy(
            QSizePolicy.Policy.Expanding, QSizePolicy.Policy.Expanding)

        # 页面 0：占位提示
        self._preview_placeholder = QLabel("选择文件查看详情")
        self._preview_placeholder.setAlignment(Qt.AlignmentFlag.AlignCenter)
        self._preview_stack.addWidget(self._preview_placeholder)

        # 页面 1：图片预览
        self._preview_image = QLabel()
        self._preview_image.setAlignment(Qt.AlignmentFlag.AlignCenter)
        self._preview_stack.addWidget(self._preview_image)

        # 页面 2：文本 / 代码预览
        self._preview_text = QTextEdit()
        self._preview_text.setReadOnly(True)
        self._preview_stack.addWidget(self._preview_text)

        # 页面 3：PDF 预览
        self._preview_pdf = QScrollArea()
        self._preview_pdf.setWidgetResizable(True)
        self._preview_pdf_container = QWidget()
        self._preview_pdf_container.setStyleSheet("background: transparent;")
        self._preview_pdf_layout = QVBoxLayout(self._preview_pdf_container)
        self._preview_pdf_layout.setContentsMargins(4, 4, 4, 4)
        self._preview_pdf_layout.setSpacing(4)
        self._preview_pdf_layout.addStretch()
        self._preview_pdf.setWidget(self._preview_pdf_container)
        self._preview_stack.addWidget(self._preview_pdf)

        self._preview_layout.addWidget(self._preview_stack)

        # 文件信息区域
        self._preview_info = QLabel()
        self._preview_info.setWordWrap(True)
        self._preview_info.setStyleSheet(
            "font-size: 11px; border: none; background: transparent;")
        self._preview_layout.addWidget(self._preview_info)

        self._preview_layout.addStretch()

        # 快捷操作按钮
        btn_layout = QHBoxLayout()
        self._preview_open_btn = QPushButton("打开")
        self._preview_open_btn.setVisible(False)
        btn_layout.addWidget(self._preview_open_btn)

        self._preview_folder_btn = QPushButton("打开目录")
        self._preview_folder_btn.setVisible(False)
        btn_layout.addWidget(self._preview_folder_btn)

        btn_layout.addStretch()
        self._preview_layout.addLayout(btn_layout)

        return panel

    def _show_preview(self, file_id: int):
        """显示文件预览（支持图片/文本/代码/PDF）"""
        record = self._file_model.get_record_by_id(file_id)
        if record is None:
            record = self.file_dao.get_by_id(file_id)
        if not record:
            self._clear_preview()
            return

        # 展开预览面板
        if not self._preview_panel.isVisible():
            self._preview_panel.setVisible(True)
            self._splitter.setSizes([180, 580, 230])

        file_path = record.get('file_path', '')
        file_name = record.get('file_name', '')
        file_type = record.get('file_type', 'other')
        file_size = record.get('file_size', 0)
        mtime = record.get('modify_time', '')
        ext = record.get('file_extension', '').lower()

        self._preview_title.setText(f"📝 {file_name[:20]}")
        self._preview_title.setWordWrap(True)

        # 基本信息
        type_name = FILE_TYPE_NAMES.get(file_type, file_type)
        info_lines = [
            f"类型: {type_name} ({ext})",
            f"大小: {format_size(file_size)}",
            f"修改: {mtime or '-'}",
            f"路径: {truncate_path(file_path, 40)}",
        ]

        # 尝试加载元数据
        try:
            meta = self.meta_dao.get_by_file_id(file_id)
            if meta:
                if meta.get('width') and meta.get('height'):
                    info_lines.append(f"尺寸: {meta['width']}×{meta['height']}")
                if meta.get('photo_taken_time'):
                    info_lines.append(f"拍摄: {meta['photo_taken_time']}")
                if meta.get('camera_model'):
                    info_lines.append(f"相机: {meta['camera_model']}")
                if meta.get('pdf_title'):
                    info_lines.append(f"PDF标题: {meta['pdf_title'][:30]}")
                if meta.get('pdf_pages'):
                    info_lines.append(f"PDF页数: {meta['pdf_pages']}")
        except Exception:
            pass

        self._preview_info.setText("\n".join(info_lines))

        # 根据文件类型路由到对应预览方法
        if file_type == 'image' and os.path.exists(file_path):
            self._preview_image_file(file_path)
        elif ext in self.PDF_EXTS and os.path.exists(file_path):
            self._preview_pdf_file(file_path, file_name)
        elif ext in self.DOCX_EXTS and os.path.exists(file_path):
            self._preview_office_file(file_path, file_name, ext)
        elif ext in self.DOC_EXTS and os.path.exists(file_path):
            self._preview_legacy_doc(file_path, file_name, ext)
        elif self._is_text_file(ext) and os.path.exists(file_path):
            self._preview_text_file(file_path, file_name, ext)
        else:
            # 不支持预览的类型：显示图标占位
            self._preview_stack.setCurrentIndex(0)
            self._preview_placeholder.setText(f"{get_file_icon(file_type)}\n{type_name}")
            self._preview_placeholder.setStyleSheet(
                "font-size: 14px; background: transparent;")
            self._preview_panel.setMinimumWidth(230)
            self._splitter.setSizes([180, 580, 230])

        # 快捷按钮
        self._preview_open_btn.setVisible(True)
        self._preview_folder_btn.setVisible(True)
        self._current_preview_path = file_path
        self._current_preview_id = file_id

        try:
            self._preview_open_btn.clicked.disconnect()
            self._preview_folder_btn.clicked.disconnect()
        except TypeError:
            pass
        self._preview_open_btn.clicked.connect(self._on_open_file_clicked)
        self._preview_folder_btn.clicked.connect(self._on_open_folder_clicked)

    # ── 各类型预览子方法 ──

    def _preview_image_file(self, file_path):
        """图片文件预览"""
        self._preview_stack.setCurrentIndex(1)
        self._preview_panel.setMinimumWidth(230)
        self._splitter.setSizes([180, 580, 230])
        try:
            pixmap = QPixmap(file_path)
            if not pixmap.isNull():
                # 高 DPI 适配：按屏幕像素比增大渲染画布，避免缩放模糊
                dpr = self.devicePixelRatio()
                base_w, base_h = 280, 200
                scaled = pixmap.scaled(
                    QSize(int(base_w * dpr), int(base_h * dpr)),
                    Qt.AspectRatioMode.KeepAspectRatio,
                    Qt.TransformationMode.SmoothTransformation)
                scaled.setDevicePixelRatio(dpr)
                self._preview_image.setPixmap(scaled)
            else:
                self._preview_image.setText("无法加载图片")
        except Exception:
            self._preview_image.setText("无法加载图片")

    def _preview_text_file(self, file_path, file_name, ext):
        """文本文件预览（TXT/CSV/LOG/MD 等）"""
        content, encoding = self._read_file_content(file_path)
        if content is None:
            self._preview_stack.setCurrentIndex(0)
            self._preview_placeholder.setText("无法读取文件内容")
            self._preview_placeholder.setStyleSheet(
                "font-size: 11px; background: transparent;")
            return

        self._preview_stack.setCurrentIndex(2)
        self._preview_panel.setMinimumWidth(380)
        self._splitter.setSizes([180, 430, 380])

        # 代码文件启用语法高亮，纯文本文件直接显示
        if ext in self.CODE_EXTS:
            self._preview_code_content(file_path, file_name, content)
        else:
            self._preview_text.setPlainText(content)

        self._preview_title.setText(f"📄 {file_name[:18]} ({encoding})")

    def _preview_code_content(self, file_path, file_name, content):
        """代码文件语法高亮"""
        try:
            from pygments import highlight
            from pygments.lexers import get_lexer_for_filename, TextLexer
            from pygments.formatters import HtmlFormatter

            try:
                lexer = get_lexer_for_filename(file_name)
            except Exception:
                lexer = TextLexer()

            formatter = HtmlFormatter(
                style='monokai', full=False, linenos=False,
                noclasses=True, nowrap=False,
            )
            highlighted = highlight(content, lexer, formatter)
            self._preview_text.setHtml(highlighted)
        except Exception:
            self._preview_text.setPlainText(content)

    def _preview_office_file(self, file_path, file_name, ext):
        """Office 文档预览（.docx/.pptx 文本提取）"""
        self._preview_stack.setCurrentIndex(2)
        self._preview_panel.setMinimumWidth(380)
        self._splitter.setSizes([180, 430, 380])

        try:
            if ext == '.docx':
                from docx import Document
                doc = Document(file_path)
                try:
                    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
                except AttributeError as attr_err:
                    # python-docx 已知 bug：含内联图片的文档遍历 paragraphs
                    # 时会触发 'list' object has no attribute 'rId'。
                    # 回退到 lxml 直接提取所有 <w:t> 文本节点。
                    if 'rId' in str(attr_err):
                        logger.warning(
                            f"python-docx paragraphs 遍历失败（rId 错误），回退 XML 提取: {attr_err}")
                        import zipfile
                        from lxml import etree
                        nsmap = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}
                        texts = []
                        with zipfile.ZipFile(file_path) as zf:
                            with zf.open('word/document.xml') as xml_f:
                                tree = etree.parse(xml_f)
                                for t_node in tree.iterfind('.//w:t', nsmap):
                                    if t_node.text:
                                        texts.append(t_node.text)
                        paragraphs = texts[:400]  # 最多 400 个文本片段
                    else:
                        raise
                content = "\n".join(paragraphs[:200])  # 最多 200 段
                icon = "📄"
                title_suffix = "Word 文档"
            elif ext == '.pptx':
                from pptx import Presentation
                prs = Presentation(file_path)
                slides_text = []
                for i, slide in enumerate(prs.slides[:10]):  # 最多 10 张幻灯片
                    slide_texts = []
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                t = para.text.strip()
                                if t:
                                    slide_texts.append(t)
                    if slide_texts:
                        slides_text.append(f"── 幻灯片 {i+1} ──\n" + "\n".join(slide_texts))
                content = "\n\n".join(slides_text)
                icon = "📊"
                title_suffix = "PPT 演示文稿"
            else:
                content = "不支持的 Office 格式"
                icon = "📄"
                title_suffix = "文档"

            if not content:
                content = "(文档无文本内容)"

            self._preview_text.setPlainText(content)
            self._preview_title.setText(f"{icon} {file_name[:18]} ({title_suffix})")

        except ImportError as e:
            missing = str(e).split()[-1] if str(e).split() else '未知'
            self._preview_text.setPlainText(f"预览需安装 {missing}\npip install python-docx python-pptx")
        except Exception as e:
            self._preview_text.setPlainText(f"文档预览失败: {e}")

    def _preview_legacy_doc(self, file_path, file_name, ext):
        """旧版 .doc/.ppt 文件预览（二进制格式，尝试提取可读文本）"""
        self._preview_stack.setCurrentIndex(2)
        self._preview_panel.setMinimumWidth(380)
        self._splitter.setSizes([180, 430, 380])

        # 尝试从二进制中提取可读文本段（UTF-16 LE 编码的文本可能嵌入在 .doc 中）
        try:
            with open(file_path, 'rb') as f:
                data = f.read(100000)

            # 尝试 UTF-16 LE 解码（旧 .doc 常用编码）
            try:
                text = data.decode('utf-16-le', errors='ignore')
                # 过滤掉控制字符，保留可打印文本
                import re
                clean = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', '', text)
                # 提取连续的中文/英文段落
                parts = [p.strip() for p in clean.split('\x00') if len(p.strip()) > 2]
                content = "\n".join(parts[:100]) if parts else None
            except Exception:
                content = None

            if content and len(content) > 20:
                self._preview_text.setPlainText(content[:50000])
                self._preview_title.setText(f"📄 {file_name[:18]} (旧版文档, 文本提取)")
            else:
                self._preview_text.setPlainText(
                    "旧版 .doc 格式无法直接预览\n"
                    "请使用「打开」按钮用 Word 查看\n\n"
                    "提示：可在 Word 中另存为 .docx 以支持预览"
                )
                self._preview_title.setText(f"📄 {file_name[:18]} (旧版 .doc)")
        except Exception:
            self._preview_text.setPlainText("旧版文档读取失败，请使用「打开」按钮查看")

    def _preview_pdf_file(self, file_path, file_name):
        """PDF 文件预览（渲染页面为图片）"""
        # 清空旧页面
        while self._preview_pdf_layout.count() > 1:
            item = self._preview_pdf_layout.takeAt(0)
            if item.widget():
                item.widget().deleteLater()

        self._preview_stack.setCurrentIndex(3)
        self._preview_panel.setMinimumWidth(380)
        self._splitter.setSizes([180, 430, 380])

        try:
            import fitz
            doc = fitz.open(file_path)
            total_pages = doc.page_count
            max_pages = min(total_pages, 20)
            # 高 DPI：按屏幕像素比放大渲染，确保 PDF 文字清晰
            dpr = self.devicePixelRatio()
            preview_width = int(350 * dpr)

            for i in range(max_pages):
                page = doc[i]
                zoom = preview_width / page.rect.width
                mat = fitz.Matrix(zoom, zoom)
                pix = page.get_pixmap(matrix=mat)
                img_data = pix.tobytes("png")

                page_pixmap = QPixmap()
                page_pixmap.loadFromData(img_data)
                page_pixmap.setDevicePixelRatio(dpr)

                page_label = QLabel()
                page_label.setPixmap(page_pixmap)
                page_label.setAlignment(Qt.AlignmentFlag.AlignCenter)

                page_num_label = QLabel(f"第 {i+1}/{total_pages} 页")
                page_subtle = "#a6adc8" if getattr(self, '_theme', 'dark') == 'dark' else "#6c7086"
                page_num_label.setStyleSheet(
                    f"color: {page_subtle}; font-size: 10px; font-weight: bold; "
                    "border: none; background: transparent;")
                page_num_label.setAlignment(Qt.AlignmentFlag.AlignCenter)

                insert_pos = self._preview_pdf_layout.count() - 1
                self._preview_pdf_layout.insertWidget(insert_pos, page_label)
                self._preview_pdf_layout.insertWidget(insert_pos + 1, page_num_label)

            doc.close()

            if total_pages > max_pages:
                more_label = QLabel(f"... 共 {total_pages} 页，仅预览前 {max_pages} 页")
                page_muted = "#585b70" if getattr(self, '_theme', 'dark') == 'dark' else "#9ca0b0"
                more_label.setStyleSheet(
                    f"color: {page_muted}; font-size: 10px; border: none; background: transparent;")
                more_label.setAlignment(Qt.AlignmentFlag.AlignCenter)
                insert_pos = self._preview_pdf_layout.count() - 1
                self._preview_pdf_layout.insertWidget(insert_pos, more_label)

            self._preview_pdf.verticalScrollBar().setValue(0)

        except ImportError:
            self._preview_stack.setCurrentIndex(0)
            self._preview_placeholder.setText("PDF 预览需 PyMuPDF\npip install PyMuPDF")
            self._preview_placeholder.setStyleSheet(
                "font-size: 11px; background: transparent;")
        except Exception as e:
            # 渲染失败则降级为 PyPDF2 文本提取
            self._preview_stack.setCurrentIndex(2)
            self._preview_panel.setMinimumWidth(380)
            self._splitter.setSizes([180, 430, 380])
            try:
                from PyPDF2 import PdfReader
                reader = PdfReader(file_path)
                text_parts = []
                for i, page in enumerate(reader.pages[:5]):
                    text = page.extract_text()
                    if text:
                        text_parts.append(f"--- 第 {i+1}/{len(reader.pages)} 页 ---\n{text}")
                if text_parts:
                    self._preview_text.setPlainText("\n\n".join(text_parts))
                    self._preview_title.setText(f"📕 {file_name[:18]} (文本提取)")
                else:
                    self._preview_text.setPlainText("PDF 无文本内容或内容无法提取")
            except Exception:
                self._preview_text.setPlainText(f"PDF 预览失败: {e}")

    def _clear_preview(self):
        """清空并折叠预览面板"""
        self._current_preview_path = None
        self._current_preview_id = None
        self._preview_title.setText("文件预览")
        self._preview_stack.setCurrentIndex(0)
        self._preview_placeholder.setText("选择文件查看详情")
        self._preview_placeholder.setStyleSheet(
            "font-size: 11px; background: transparent;")
        self._preview_image.clear()
        self._preview_text.clear()
        # 清空 PDF 页面
        while self._preview_pdf_layout.count() > 1:
            item = self._preview_pdf_layout.takeAt(0)
            if item.widget():
                item.widget().deleteLater()
        self._preview_info.setText("")
        self._preview_open_btn.setVisible(False)
        self._preview_folder_btn.setVisible(False)
        self._preview_panel.setMinimumWidth(230)
        # 折叠面板
        if self._preview_panel.isVisible():
            self._preview_panel.setVisible(False)
            self._splitter.setSizes([180, 820, 0])

    def apply_theme(self, tn: str):
        """主题切换时更新预览面板所有控件的样式"""
        self._theme = tn
        is_dark = tn == 'dark'

        # 颜色方案
        panel_bg = "#1e1e2e" if is_dark else "#eff1f5"
        panel_border = "#313244" if is_dark else "#ccd0da"
        content_bg = "#181825" if is_dark else "#e6e9ef"
        text_color = "#cdd6f4" if is_dark else "#4c4f69"
        muted_color = "#585b70" if is_dark else "#9ca0b0"
        subtle_color = "#a6adc8" if is_dark else "#6c7086"
        accent_color = "#cba6f7" if is_dark else "#8839ef"
        btn_bg = "#313244" if is_dark else "#ccd0da"
        btn_hover = "#45475a" if is_dark else "#bcc0cc"
        btn_border = "#45475a" if is_dark else "#bcc0cc"
        scroll_bg = "#181825" if is_dark else "#e6e9ef"
        scroll_handle = "#45475a" if is_dark else "#ccd0da"

        # ── 预览面板外框 ──
        if hasattr(self, '_preview_panel'):
            self._preview_panel.setStyleSheet(
                f"QFrame {{ background-color: {panel_bg}; "
                f"border: 1px solid {panel_border}; border-radius: 6px; }}")

        # ── 标题 ──
        if hasattr(self, '_preview_title'):
            self._preview_title.setStyleSheet(
                f"font-weight: bold; color: {accent_color}; font-size: 12px;"
                "border: none; background: transparent;")

        # ── 预览栈背景 ──
        if hasattr(self, '_preview_stack'):
            self._preview_stack.setStyleSheet(
                f"background-color: {content_bg}; border-radius: 4px; border: none;")

        # ── 占位/图片标签 ──
        if hasattr(self, '_preview_placeholder'):
            self._preview_placeholder.setStyleSheet(
                f"color: {muted_color}; font-size: 11px; background: transparent;")
        if hasattr(self, '_preview_image'):
            self._preview_image.setStyleSheet(
                f"color: {muted_color}; background: transparent;")

        # ── 文本预览 (QTextEdit) ──
        if hasattr(self, '_preview_text'):
            self._preview_text.setStyleSheet(
                f"QTextEdit {{"
                f"  color: {text_color};"
                f"  background-color: {content_bg};"
                "  border: none;"
                "  font-family: 'Consolas', 'Courier New', monospace;"
                "  font-size: 11px;"
                f"}}"
                f"QScrollBar:vertical {{ background: {scroll_bg}; width: 6px; }}"
                f"QScrollBar::handle:vertical {{ background: {scroll_handle}; border-radius: 3px; }}"
                "QScrollBar::add-line:vertical, QScrollBar::sub-line:vertical { height: 0; }")

        # ── PDF 滚动区域 ──
        if hasattr(self, '_preview_pdf'):
            self._preview_pdf.setStyleSheet(
                f"QScrollArea {{ background-color: {content_bg}; border: none; }}"
                f"QScrollBar:vertical {{ background: {scroll_bg}; width: 6px; }}"
                f"QScrollBar::handle:vertical {{ background: {scroll_handle}; border-radius: 3px; }}"
                "QScrollBar::add-line:vertical, QScrollBar::sub-line:vertical { height: 0; }")

        # ── 文件信息 ──
        if hasattr(self, '_preview_info'):
            self._preview_info.setStyleSheet(
                f"color: {text_color}; font-size: 11px; border: none; background: transparent;")

        # ── 快捷按钮 ──
        if hasattr(self, '_preview_open_btn'):
            self._preview_open_btn.setStyleSheet(
                f"QPushButton {{ background-color: {btn_bg}; color: {text_color}; "
                f"border: 1px solid {btn_border}; border-radius: 4px; "
                "padding: 4px 8px; font-size: 11px; }"
                f"QPushButton:hover {{ background-color: {btn_hover}; }}")
        if hasattr(self, '_preview_folder_btn'):
            self._preview_folder_btn.setStyleSheet(
                f"QPushButton {{ background-color: {btn_bg}; color: {text_color}; "
                f"border: 1px solid {btn_border}; border-radius: 4px; "
                "padding: 4px 8px; font-size: 11px; }"
                f"QPushButton:hover {{ background-color: {btn_hover}; }}")

        # ── 刷新当前预览（如有）以更新动态样式 ──
        if (hasattr(self, '_current_preview_path')
                and self._current_preview_path
                and hasattr(self, '_preview_stack')):
            idx = self._preview_stack.currentIndex()
            if idx == 0:
                # 占位页面：更新颜色
                if hasattr(self, '_preview_placeholder'):
                    self._preview_placeholder.setStyleSheet(
                        f"color: {muted_color}; font-size: 11px; background: transparent;")

    def _on_open_file_clicked(self):
        """打开文件按钮点击"""
        if hasattr(self, '_current_preview_path') and self._current_preview_path:
            self._safe_open_file(self._current_preview_path, file_id=getattr(self, '_current_preview_id', None))

    def _on_open_folder_clicked(self):
        """打开文件夹按钮点击"""
        if hasattr(self, '_current_preview_path') and self._current_preview_path:
            self._safe_open_folder(self._current_preview_path, file_id=getattr(self, '_current_preview_id', None))
